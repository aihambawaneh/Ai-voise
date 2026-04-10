import os
import json
import asyncio
import time
import google.generativeai as genai
import edge_tts
from pptx import Presentation
from dotenv import load_dotenv

# Load env in brain too
load_dotenv()

# Support multiple keys for load balancing/quota
api_keys = os.getenv("GEMINI_API_KEYS", "").split(",")
if not api_keys or not api_keys[0]:
    # Fallback to single key if only one is provided
    api_keys = [os.getenv("GEMINI_API_KEY")]

current_key_header = 0
genai.configure(api_key=api_keys[0])

def get_best_model():
    try:
        models = [m.name for m in genai.list_models()]
        print(f"📦 Available models: {models}")
        
        for m in genai.list_models():
            if 'generateContent' in m.supported_generation_methods:
                if '1.5-flash' in m.name:
                    return m.name
        for m in genai.list_models():
            if 'generateContent' in m.supported_generation_methods:
                return m.name
    except Exception as e:
        print(f"❌ Error listing models: {e}")
    return 'gemini-1.5-flash'

class AI_Brain:
    def __init__(self):
        self.model_name = get_best_model()
        print(f"USING MODEL: {self.model_name}")
        self.model = genai.GenerativeModel(
            self.model_name,
            generation_config={"response_mime_type": "application/json"}
        )

    async def generate_audio(self, text, output_path, voice):
        communicate = edge_tts.Communicate(text, voice)
        await communicate.save(output_path)

    async def generate_lecture_data(self, topic, language, level='student'):
        # Dynamic style instruction based on Level
        num_slides = "4-6"
        if level == 'kid':
            num_slides = "3-5"
            style_instruction = "شرح للأطفال الصغار (5-8 سنوات). خاطبهم بـ 'يا أبطال' أو 'يا حلوين'. استخدم أسلوب القصص والتبسيط الشديد. اجعل الكلام مشوقاً ومرحاً جداً وكأنك في برنامج أطفال."
        elif level == 'pro':
            num_slides = "8-12"
            style_instruction = "محاضرة جامعية احترافية للأكاديميين. خاطبهم بأسلوب رسمي مثل 'أعزائي المختصين' أو 'الزملاء الكرام' أو ابدأ مباشرة بجدية. استخدم مصطلحات تقنية عميقة ومحتوى مكثف جداً. لا تستخدم أي كلمات طفولية أو تبسيط مخل."
        else:
            style_instruction = "شرح مدرسي متوازن لطلاب المتوسط/الثانوي. خاطبهم كشباب ناضجين بأسلوب 'يا شباب' أو 'يا مبدعين'. استخدم توازن بين البساطة والتفاصيل العلمية."

        # Dialect specific refinements
        dialect_instruction = ""
        if language == 'Arabic Saudi':
            dialect_instruction = """
            التعليمات للهجة السعودية (حاسمة):
            - لا تكتب باللغة العربية الفصحى أبداً. اكتب كما يتحدث الناس في شوارع الرياض وجدة.
            - استخدم كلمات مثل: 'يا هلا والله'، 'بشرونا عنكم'، 'وش السالفة'، 'قاعدين نشرح'، 'شوفوا معي'، 'لازم ننتبه'، 'عشان كذا'.
            - استبدل 'لماذا' بـ 'ليه'، و'كيف' بـ 'شلون' أو 'كيف'، و'هنا' بـ 'هنا' أو 'هنيا'.
            - مثال للأسلوب المطلوب: 'يا هلا فيكم اليوم في درسنا الجديد.. اليوم نبي نسولف عن علم الفلك وشلون النجوم تطلع في السماء.. شوفوا معي هالنقطة ترى مرة مهمة'.
            - اجعل الأسلوب 'سالفة' وليس 'إلقاء كتابي'.
            """
        elif language == 'Arabic Fusha':
            dialect_instruction = "يجب الالتزام باللغة العربية الفصحى الحديثة (MSA) التزاماً تاماً، وبأسلوب فخم ورصين."

        prompt = f"""
        الهدف: إنشاء محتوى تعليمي احترافي لدرس عن: '{topic}'.
        المستوى المطلوب: {level} ({style_instruction}).
        اللغة/اللهجة المطلوبة: {language}.
        {dialect_instruction}
        
        عدد الشرائح المطلوبة: {num_slides}.
        
        المطلوب منك هو إرجاع مخرجات بصيغة JSON حصراً، وتحتوي على مصفوفة من الكائنات:
        كل كائن يجب أن يحتوي على:
        - "title": عنوان شريحة قصير ومبهر.
        - "points": مصفوفة تحتوي على 3-5 نقاط (مختصرة جداً - سطر واحد لكل نقطة).
        - "explanation": نص الشرح الصوتي للدرس. هنا يجب أن يكون الشرح (مكثفاً جداً، طويلاً، ومليئاً بالمعلومات التقنية العميقة للمستوى Pro).
        
        تذكر: الكلام الذي يظهر في الشريحة (points) يجب أن يكون رؤوس أقلام فقط، أما العلم الحقيقي والمعلومات الدقيقة تكون في (explanation).
        
        Output format (STRICT JSON):
        [
          {{
            "title": "...",
            "points": ["...", "...", "..."],
            "explanation": "..."
          }}
        ]
        """

        # Handle Quota/Rate Limits with a simple retry
        max_retries = 3
        for attempt in range(max_retries):
            try:
                response = self.model.generate_content(prompt)
                text_content = response.text.strip()
                
                try:
                    return json.loads(text_content)
                except json.JSONDecodeError:
                    import re
                    match = re.search(r'\[.*\]', text_content, re.DOTALL)
                    if match:
                        return json.loads(match.group())
                    raise ValueError("Could not parse JSON. Raw: " + text_content)
                    
            except Exception as e:
                if "429" in str(e) and attempt < max_retries - 1:
                    wait_time = (attempt + 1) * 2
                    
                    # Switch to next key if available
                    global current_key_header
                    current_key_header = (current_key_header + 1) % len(api_keys)
                    new_key = api_keys[current_key_header]
                    genai.configure(api_key=new_key)
                    self.model = genai.GenerativeModel(self.model_name, generation_config={"response_mime_type": "application/json"})
                    
                    print(f"⚠️ Quota hit. Switching to key {current_key_header + 1} and retrying in {wait_time}s...")
                    import time
                    time.sleep(wait_time)
                    continue
                raise e

    def create_pptx(self, slides_data, output_path):
        prs = Presentation()
        for i, slide_info in enumerate(slides_data):
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = slide_info.get('title', f"Slide {i+1}")
            tf = slide.placeholders[1].text_frame
            tf.text = ""
            for point in slide_info.get('points', []):
                p = tf.add_paragraph()
                p.text = point
        prs.save(output_path)
